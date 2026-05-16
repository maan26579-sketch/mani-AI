#!/usr/bin/env python3
"""
generate_presentation.py

Generates Ukraine-Russia Conflict Presentation (.pptx) from embedded slide data and images.

Usage:
  pip install python-pptx pillow requests matplotlib
  python generate_presentation.py

Output:
  Ukraine-Russia_Conflict_Presentation.pptx

This script was committed to the repo so you (or teammates) can run it locally to produce the binary .pptx file.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import requests
from io import BytesIO
from PIL import Image
import matplotlib.pyplot as plt
import os

OUTPUT = 'Ukraine-Russia_Conflict_Presentation.pptx'

# Slides data: title, bullets (list), speaker notes, image_url (optional)
SLIDES = [
    ("Ukraine–Russia Conflict — Overview & Current Issues",
     ["Class presentation • [Your class name] • Date: [DD MMM YYYY]"],
     "One-sentence intro: We are [your group name]; today we’ll present a concise overview of the Ukraine–Russia conflict: background, timeline, impacts, international response, and outlook.",
     "https://upload.wikimedia.org/wikipedia/commons/9/90/Political_Map_of_Ukraine.svg.png"),

    ("Background / Historical Context",
     ["Ukraine: Soviet republic → independent in 1991",
      "Complex historical ties to Russia; regional and linguistic diversity",
      "Strategic importance: Black Sea ports, agriculture, energy transit routes"],
     "Summarize key historical drivers (Soviet history, Crimea’s 1954 transfer, cultural/regional splits).",
     "https://upload.wikimedia.org/wikipedia/commons/9/90/Political_Map_of_Ukraine.svg.png"),

    ("Key Events through 2014",
     ["2004 Orange Revolution; 2013–14 Euromaidan protests",
      "2014: Russia annexes Crimea; armed conflict in Donbas begins",
      "Minsk I/II attempted ceasefires; periodic violations and frozen lines"],
     "Explain how Euromaidan led to regime change and Russia’s reaction in Crimea/Donbas.",
     "https://upload.wikimedia.org/wikipedia/commons/8/81/Ukraine_map_%28disputed_territory%29.png"),

    ("2022 Full‑Scale Invasion: Overview",
     ["24 Feb 2022: Russia launches a large-scale military invasion",
      "Initial multi-axis offensive toward Kyiv, advances in east and south",
      "Subsequent phases: Kyiv withdrawal, focus on east/south, Ukrainian counteroffensives"],
     "Concise operational summary; note shifting frontlines and phases.",
     "https://upload.wikimedia.org/wikipedia/commons/9/96/UKR_Territorial_control_2022-11-12.png"),

    ("Military & Tactical Aspects",
     ["Combined-arms warfare: ground forces, artillery, air, missile strikes, drones",
      "Logistics, mobilization, and attrition shape operational tempo",
      "External military aid, training, and ISR influence capabilities"],
     "Explain modern elements (artillery, drones, air defenses) and logistics impacts.",
     None),

    ("Humanitarian Impact",
     ["Large civilian toll: deaths and injuries; hospitals, homes, schools damaged",
      "Refugees abroad and millions internally displaced",
      "Urgent humanitarian needs: shelter, food, medical care, infrastructure repair"],
     "Includes figures: Total civilian casualties: 58,930 (13,883 killed and 43,352 injured). Refugees abroad: 6.7M. IDPs: 6.5M. (Sources: OHCHR, UNHCR, OCHA)",
     None),

    ("Economic Consequences & Sanctions",
     ["Broad sanctions: banking, SWIFT exclusions, asset freezes, export controls",
      "Energy and food market disruption (grain exports, fuel prices)",
      "Long-term reconstruction costs and economic ripple effects"],
     "Explain how sanctions target finance, energy, and tech; note global impacts.",
     None),

    ("International Response",
     ["Diplomatic measures: UN debates, sanctions packages, political isolation",
      "Military assistance: arms, training, intelligence sharing for Ukraine",
      "Humanitarian assistance and multi-lateral coordination (UN, EU, NGOs)"],
     "Distinguish diplomatic tools from military assistance; mention major contributors.",
     None),

    ("Information Warfare, Legal Issues, & Cybersecurity",
     ["Disinformation campaigns and social media narratives",
      "Cyberattacks on government, infrastructure, and media organizations",
      "War crimes investigations: ICC, national investigations; evidence & prosecution challenges"],
     "Give one example of information warfare tactics and the effects on perception.",
     None),

    ("Conclusion, Sources & Q&A",
     ["Key takeaways: complex historical causes, severe humanitarian toll, global impact, uncertain outcome",
      "Next steps: verify latest data, support humanitarian response, follow reputable sources",
      "Sources: OHCHR, UNHCR, OCHA, ICRC, ReliefWeb, major news outlets"],
     "Summarize and invite questions. Point audience to sources for links.",
     None)
]


def add_title_slide(prs, title, subtitle, notes, image_url=None):
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    # Add background image if available
    if image_url:
        try:
            img_bytes = fetch_image_bytes(image_url)
            slide.shapes.add_picture(img_bytes, 0, 0, width=prs.slide_width)
        except Exception:
            pass
    # Add title textbox
    left = Inches(0.5)
    top = Inches(0.7)
    width = prs.slide_width - Inches(1)
    height = Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(10, 57, 100)

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(60, 60, 60)

    slide.notes_slide.notes_text_frame.text = notes


def add_content_slide(prs, title, bullets, notes, image_url=None):
    slide_layout = prs.slide_layouts[1]  # title and content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    # content placeholder
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.font.size = Pt(18)
    # Add image to right if provided
    if image_url:
        try:
            img_bytes = fetch_image_bytes(image_url)
            pic_left = prs.slide_width - Inches(4.0)
            pic_top = Inches(1.5)
            slide.shapes.add_picture(img_bytes, pic_left, pic_top, width=Inches(3.5))
        except Exception:
            pass
    # Add notes
    slide.notes_slide.notes_text_frame.text = notes


def fetch_image_bytes(url):
    r = requests.get(url, timeout=10)
    r.raise_for_status()
    return BytesIO(r.content)


def create_chart_image():
    # Creates a small bar chart comparing Refugees vs IDPs
    labels = ['Refugees', 'IDPs']
    values = [6.7, 6.5]
    fig, ax = plt.subplots(figsize=(4,2.5))
    bars = ax.bar(labels, values, color=['#2b7cff', '#ffa500'])
    ax.set_ylabel('Millions')
    ax.set_title('Displaced People (millions)')
    ax.set_ylim(0, max(values) + 1)
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.05, f'{val}M', ha='center', va='bottom')
    buf = BytesIO()
    plt.tight_layout()
    fig.savefig(buf, format='png', dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf


def main():
    prs = Presentation()
    # Optional: set slide size (default 10 x 7.5 inches)
    # Add title slide
    t0, bullets0, notes0, img0 = SLIDES[0]
    add_title_slide(prs, t0, bullets0[0], notes0, image_url=img0)

    # Add remaining slides
    for title, bullets, notes, img in SLIDES[1:]:
        add_content_slide(prs, title, bullets, notes, image_url=img)

    # Insert chart slide for humanitarian slide (Slide 6 index 5)
    # Find slide with title 'Humanitarian Impact'
    for slide in prs.slides:
        if slide.shapes.title and slide.shapes.title.text.strip().lower().startswith('humanitarian'):
            img_buf = create_chart_image()
            left = Inches(5.0)
            top = Inches(3.0)
            slide.shapes.add_picture(img_buf, left, top, width=Inches(4))
            break

    # Save presentation
    prs.save(OUTPUT)
    print(f'Presentation generated: {OUTPUT}')

if __name__ == '__main__':
    main()
